from __future__ import annotations

import io
import tempfile
import time
from pathlib import Path

import fitz
from docx import Document
from fastapi import UploadFile
from PIL import Image
from pptx import Presentation

from app.config import settings
from app.gpu_extractors import gpu_models
from app.schemas import DocumentMetadata, ElementType, ExtractedElement, ExtractionResponse


async def extract_upload(file: UploadFile, allowed_extensions: set[str]) -> ExtractionResponse:
    filename = file.filename or "uploaded_document"
    extension = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if extension not in allowed_extensions:
        raise ValueError(f"Unsupported document type '.{extension}'.")

    content = await file.read()
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{extension}") as tmp:
        tmp.write(content)
        path = Path(tmp.name)

    try:
        if extension == "pdf":
            return extract_pdf(path, filename, len(content))
        if extension == "docx":
            return extract_docx(path, filename, len(content))
        if extension == "pptx":
            return extract_pptx(path, filename, len(content))
        return extract_image(path, filename, len(content), extension)
    finally:
        path.unlink(missing_ok=True)


def build_response(
    *,
    filename: str,
    file_type: str,
    file_size: int,
    elements: list[ExtractedElement],
    started: float,
    page_count: int | None = None,
    slide_count: int | None = None,
    pages_processed: int = 0,
) -> ExtractionResponse:
    parts = []
    for element in elements:
        if not element.content.strip():
            continue
        if element.element_type == ElementType.IMAGE_DESCRIPTION:
            parts.append(f"[Visual: {element.content.strip()}]")
        else:
            parts.append(element.content.strip())

    return ExtractionResponse(
        metadata=DocumentMetadata(
            filename=filename,
            file_type=file_type,
            file_size_bytes=file_size,
            page_count=page_count,
            slide_count=slide_count,
        ),
        elements=elements,
        full_text="\n\n".join(parts),
        processing_time_ms=round((time.perf_counter() - started) * 1000, 2),
        pages_processed=pages_processed,
        native_count=sum(1 for e in elements if e.source == "native"),
        ocr_count=sum(1 for e in elements if e.source == "gpu_ocr"),
        visual_count=sum(1 for e in elements if e.source == "gpu_visual"),
        device_used="cuda",
    )


def extract_pdf(path: Path, filename: str, file_size: int) -> ExtractionResponse:
    started = time.perf_counter()
    elements: list[ExtractedElement] = []
    doc = fitz.open(path)
    pages_to_process = min(len(doc), settings.request_max_pages)

    for page_index in range(pages_to_process):
        page = doc[page_index]
        page_no = page_index + 1

        native_text = page.get_text("text").strip()
        if native_text:
            elements.append(
                ExtractedElement(
                    element_type=ElementType.TEXT,
                    content=f"Page {page_no} text:\n{native_text}",
                    page=page_no,
                    source="native",
                )
            )

        pixmap = page.get_pixmap(dpi=settings.pdf_render_dpi, alpha=False)
        image = Image.open(io.BytesIO(pixmap.tobytes("png")))
        visual_text = gpu_models.caption(image)
        ocr_text = gpu_models.ocr(image)

        elements.append(
            ExtractedElement(
                element_type=ElementType.IMAGE_DESCRIPTION,
                content=f"Page {page_no} visual description: {visual_text}",
                page=page_no,
                source="gpu_visual",
            )
        )
        if ocr_text and ocr_text not in native_text:
            elements.append(
                ExtractedElement(
                    element_type=ElementType.TEXT,
                    content=f"Page {page_no} OCR text:\n{ocr_text}",
                    page=page_no,
                    source="gpu_ocr",
                )
            )

    return build_response(
        filename=filename,
        file_type="pdf",
        file_size=file_size,
        elements=elements,
        started=started,
        page_count=len(doc),
        pages_processed=pages_to_process,
    )


def extract_image(path: Path, filename: str, file_size: int, file_type: str) -> ExtractionResponse:
    started = time.perf_counter()
    image = Image.open(path)
    caption = gpu_models.caption(image)
    ocr_text = gpu_models.ocr(image)

    elements = [
        ExtractedElement(
            element_type=ElementType.IMAGE_DESCRIPTION,
            content=f"Image visual description: {caption}",
            page=1,
            source="gpu_visual",
        )
    ]
    if ocr_text:
        elements.append(
            ExtractedElement(
                element_type=ElementType.TEXT,
                content=f"Image OCR text:\n{ocr_text}",
                page=1,
                source="gpu_ocr",
            )
        )

    return build_response(
        filename=filename,
        file_type=file_type,
        file_size=file_size,
        elements=elements,
        started=started,
        page_count=1,
        pages_processed=1,
    )


def extract_docx(path: Path, filename: str, file_size: int) -> ExtractionResponse:
    started = time.perf_counter()
    doc = Document(path)
    elements: list[ExtractedElement] = []

    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if paragraphs:
        elements.append(
            ExtractedElement(
                element_type=ElementType.TEXT,
                content="\n".join(paragraphs),
                source="native",
            )
        )

    for table_index, table in enumerate(doc.tables, start=1):
        rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        elements.append(
            ExtractedElement(
                element_type=ElementType.TABLE,
                content=f"Table {table_index}:\n" + "\n".join(rows),
                source="native",
            )
        )

    image_index = 0
    for rel in doc.part.rels.values():
        if "image" not in rel.reltype:
            continue
        image_index += 1
        image = Image.open(io.BytesIO(rel.target_part.blob))
        caption = gpu_models.caption(image)
        ocr_text = gpu_models.ocr(image)
        elements.append(
            ExtractedElement(
                element_type=ElementType.IMAGE_DESCRIPTION,
                content=f"Embedded image {image_index} visual description: {caption}",
                source="gpu_visual",
                metadata={"image_index": image_index},
            )
        )
        if ocr_text:
            elements.append(
                ExtractedElement(
                    element_type=ElementType.TEXT,
                    content=f"Embedded image {image_index} OCR text:\n{ocr_text}",
                    source="gpu_ocr",
                    metadata={"image_index": image_index},
                )
            )

    return build_response(
        filename=filename,
        file_type="docx",
        file_size=file_size,
        elements=elements,
        started=started,
        pages_processed=1,
    )


def extract_pptx(path: Path, filename: str, file_size: int) -> ExtractionResponse:
    started = time.perf_counter()
    prs = Presentation(path)
    elements: list[ExtractedElement] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
            if getattr(shape, "shape_type", None) == 13:
                image = Image.open(io.BytesIO(shape.image.blob))
                caption = gpu_models.caption(image)
                ocr_text = gpu_models.ocr(image)
                elements.append(
                    ExtractedElement(
                        element_type=ElementType.IMAGE_DESCRIPTION,
                        content=f"Slide {slide_index} image visual description: {caption}",
                        slide=slide_index,
                        source="gpu_visual",
                    )
                )
                if ocr_text:
                    elements.append(
                        ExtractedElement(
                            element_type=ElementType.TEXT,
                            content=f"Slide {slide_index} image OCR text:\n{ocr_text}",
                            slide=slide_index,
                            source="gpu_ocr",
                        )
                    )

        if slide_text:
            elements.insert(
                0,
                ExtractedElement(
                    element_type=ElementType.TEXT,
                    content=f"Slide {slide_index} text:\n" + "\n".join(slide_text),
                    slide=slide_index,
                    source="native",
                )
            )

    return build_response(
        filename=filename,
        file_type="pptx",
        file_size=file_size,
        elements=elements,
        started=started,
        slide_count=len(prs.slides),
        pages_processed=len(prs.slides),
    )
