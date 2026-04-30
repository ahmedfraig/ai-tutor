"""
app/services/pptx_service.py

Hybrid PPTX processor:
  - python-pptx: text frames, tables, speaker notes
  - VLM: each slide rendered as image + embedded pictures
"""
from __future__ import annotations

import io
import time
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from app.core.config import get_settings
from app.core.device import resolve_device
from app.core.logging import get_logger
from app.models.schemas import (
    DocumentMetadata,
    ElementType,
    ExtractionMode,
    ExtractionResponse,
    ExtractedElement,
)
from app.services.vlm_service import get_vlm_service

logger = get_logger(__name__)
settings = get_settings()


async def process_pptx(
    path: Path,
    filename: str,
    file_size: int,
    mode: ExtractionMode = ExtractionMode.AUTO,
    describe_visuals: bool = True,
    vlm_prompt: str = "Describe the content of this slide image in detail.",
) -> ExtractionResponse:
    t0 = time.perf_counter()
    vlm_svc = get_vlm_service()
    elements: list[ExtractedElement] = []
    device = resolve_device(settings.device)

    prs = Presentation(str(path))
    total_slides = len(prs.slides)

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []

        for shape in slide.shapes:
            # ── Text frames ──────────────────────────────────────────────────
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)

            # ── Tables ───────────────────────────────────────────────────────
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        rows.append(" | ".join(cells))
                if rows:
                    elements.append(ExtractedElement(
                        element_type=ElementType.TABLE,
                        content="\n".join(rows),
                        slide=slide_num,
                        source="native",
                    ))

            # ── Embedded pictures → VLM ───────────────────────────────────────
            if (
                describe_visuals
                and mode in (ExtractionMode.AUTO, ExtractionMode.FULL, ExtractionMode.VLM_ONLY)
                and shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            ):
                try:
                    img_bytes = shape.image.blob
                    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                    if img.width >= 80 and img.height >= 80:
                        desc = await vlm_svc.adescribe(img, vlm_prompt)
                        if desc.strip():
                            elements.append(ExtractedElement(
                                element_type=ElementType.FIGURE,
                                content=desc,
                                slide=slide_num,
                                source="vlm",
                            ))
                except Exception as exc:
                    logger.warning("pptx_image_failed", slide=slide_num, error=str(exc))

        # ── Aggregate slide text ──────────────────────────────────────────────
        if slide_texts:
            elements.append(ExtractedElement(
                element_type=ElementType.TEXT,
                content="\n".join(slide_texts),
                slide=slide_num,
                source="native",
            ))

        # ── Speaker notes ─────────────────────────────────────────────────────
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                elements.append(ExtractedElement(
                    element_type=ElementType.TEXT,
                    content=notes_text,
                    slide=slide_num,
                    source="native",
                    metadata={"is_speaker_note": True},
                ))

    elapsed_ms = (time.perf_counter() - t0) * 1000
    metadata = DocumentMetadata(
        filename=filename,
        file_type="pptx",
        file_size_bytes=file_size,
        slide_count=total_slides,
    )
    return ExtractionResponse.build(
        metadata=metadata,
        elements=elements,
        processing_time_ms=elapsed_ms,
        device_used=device,
        pages_processed=total_slides,
    )
